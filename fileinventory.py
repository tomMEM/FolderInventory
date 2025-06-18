"""
File Inventory Dashboard
Version: 1.0
Author: 
Description: A Gradio-based dashboard for managing file inventories with search, filtering, and note-taking capabilities.
Last Updated: 2024-06-14
"""

# ==============================================================================
# FINAL, STABLE, AND WORKING SCRIPT (VERSION 6 - ICON FIX)
# ==============================================================================
import os
import sys
import time
import threading
import datetime
import openpyxl
from openpyxl.utils import get_column_letter
import platform
import subprocess
import gradio as gr
import pandas as pd
import webbrowser
import traceback
import logging

# --- PyInstaller Windowed Mode Fix ---
if sys.stderr is None:
    class DummyStream:
        def isatty(self): return False
        def write(self, msg): pass
        def flush(self): pass
    sys.stderr = DummyStream()
    sys.stdout = DummyStream()

# --- Proxy Setup (omitted for brevity, assume it's correct) ---
PROXY_URL = os.environ.get("HTTP_PROXY_URL", "http://localhost:4321")
if PROXY_URL:
    os.environ['http_proxy'] = PROXY_URL
    os.environ['https_proxy'] = PROXY_URL
    # print(f"Using HTTP_PROXY: {os.environ.get('http_proxy')}")
    # print(f"Using HTTPS_PROXY: {os.environ.get('https_proxy')}")
else:
    pass
    # print("No HTTP_PROXY_URL environment variable set, not using proxy for HTTP/HTTPS.")
current_no_proxy = os.environ.get('NO_PROXY', '')
additional_no_proxy_hosts = ['localhost', '127.0.0.1', '0.0.0.0']
new_no_proxy_parts = [host for host in current_no_proxy.split(',') if host.strip()]
for host in additional_no_proxy_hosts:
    if host not in new_no_proxy_parts:
        new_no_proxy_parts.append(host)
os.environ['NO_PROXY'] = ','.join(new_no_proxy_parts)
# print(f"Using NO_PROXY: {os.environ.get('NO_PROXY')}")


# --- Configuration ---
VERSION = "1.0.0"
BUILD_DATE = "2025-06-14"
START_FOLDER = r"\yours\folder"
INVENTORY_FILENAME = "inventory.xlsx"

# Topic keywords configuration
TOPIC_KEYWORDS = {
    "PET": {
        "all_required": ["pet"],
        "any_of_these": ["scan", "imaging", "tracer"]
    },
    "DID": {
        "all_required": ["did"],
        "any_of_these": ["dementia", "cognitive", "decline"]
    },
    "AD": {
        "all_required": ["alzheimer"],
        "any_of_these": ["disease", "dementia", "ad"]
    }
}

# File type configurations
TEXT_BASED_EXTENSIONS = [
    '.docx', '.pptx', '.txt', 
    '.py', '.r', '.md'
]
SPREADSHEET_EXTENSIONS = [
    '.xlsx', '.csv', '.prism'
]

# Column definitions
FIELDNAMES = [
    'Folder Path', 'File Name', 'Extension', 
    'Size (Bytes)', 'Last Modified', 'Full Path', 
    'Content Hint', 'Identified Topics (DOCX)', 
    'Status', 'Manual_Notes'
]

DISPLAY_COLUMNS = [
    'Action', 'File Name', 'Status', 
    'Last Modified', 'Identified Topics (DOCX)', 
    'Content Hint', 'Manual_Notes', 'Full Path'
]

STATE_COLUMNS = FIELDNAMES + ['Action'] # Columns for the full_df_state

# Error handling configuration
MAX_BACKUP_FILES = 5
ERROR_WINDOW_SECONDS = 300
MAX_ERROR_COUNT = 3

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('inventory_dashboard.log'),
        logging.StreamHandler()
    ]
)


# --- Helper Functions (get_content_hint, check_docx_for_topics, load_existing_inventory, save_inventory_to_xlsx, process_folder_inventory are mostly unchanged from previous stable version, ensure they are robust) ---

def get_content_hint(filepath, extension):
    hint = "N/A"
    try:
        if extension == '.docx':
            try:
                import docx
                doc = docx.Document(filepath)
                if doc.paragraphs: hint = "First para: " + doc.paragraphs[0].text[:150] + "..."
                else: hint = "DOCX: No paragraphs found."
            except ImportError: hint = "DOCX: 'python-docx' not installed."
            except Exception: hint = "DOCX: Corrupt or unreadable."
        elif extension == '.pptx':
            try:
                import pptx
                prs = pptx.Presentation(filepath)
                if prs.slides and prs.slides[0].shapes.title: hint = "First slide title: " + prs.slides[0].shapes.title.text[:150]
                elif prs.slides: hint = "PPTX: First slide no title."
                else: hint = "PPTX: No slides."
            except ImportError: hint = "PPTX: 'python-pptx' not installed."
            except Exception: hint = "PPTX: Corrupt or unreadable."
        elif extension in TEXT_BASED_EXTENSIONS:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                lines = [line.strip() for i, line in enumerate(f) if i < 2]
                hint_text = " ".join(lines)
                hint = ("First 2 lines: " + hint_text[:200] + "...") if hint_text else f"{extension.upper()}: Empty"
        elif extension in SPREADSHEET_EXTENSIONS:
             hint = "Spreadsheet file."
    except Exception as e: hint = f"Hint Error for {os.path.basename(filepath)}: {e}"
    return hint

def check_docx_for_topics(filepath, topic_definitions):
    identified_topics = []
    try:
        # Check if file exists and is accessible first
        if not os.path.exists(filepath):
            print(f"Warning: DOCX file not found: {filepath}")
            return "N/A (File not found)"
            
        import docx
        try:
            doc = docx.Document(filepath)
        except Exception as e:
            print(f"Warning: Could not open DOCX {filepath}: {e}")
            return "N/A (Access error)"
            
        full_text = "\n".join([para.text for para in doc.paragraphs]).lower()
        if not full_text.strip(): 
            return "DOCX Empty"
            
        for topic_name, criteria in topic_definitions.items():
            all_match = all(kw.lower() in full_text for kw in criteria.get("all_required", []))
            any_match = True
            if "any_of_these" in criteria and criteria["any_of_these"]:
                any_match = any(kw.lower() in full_text for kw in criteria["any_of_these"])
            if all_match and any_match: 
                identified_topics.append(topic_name)
    except Exception as e:
        print(f"Error checking DOCX topics for {filepath}: {e}")
        return "N/A (Error reading DOCX)"
    return ", ".join(identified_topics) if identified_topics else "N/A"

def load_existing_inventory(xlsx_filepath):
    existing_data = {}
    if not os.path.exists(xlsx_filepath):
        print(f"INFO: Inventory file '{xlsx_filepath}' not found. New one will be created.")
        return existing_data
    print(f"INFO: Attempting to load inventory from '{xlsx_filepath}'...")
    try:
        df = pd.read_excel(xlsx_filepath, engine='openpyxl')
        if 'Full Path' not in df.columns:
            print(f"ERROR: 'Full Path' column missing in '{xlsx_filepath}'. Cannot process.")
            return {}
        for col_name in FIELDNAMES: # Ensure all expected data columns exist
            if col_name not in df.columns:
                print(f"INFO: Data column '{col_name}' not found in Excel. Adding it.")
                df[col_name] = '' if col_name == 'Manual_Notes' else pd.NA
        if 'Manual_Notes' in df.columns: # Ensure notes are string
            df['Manual_Notes'] = df['Manual_Notes'].fillna('').astype(str)
        df.set_index('Full Path', inplace=True)
        existing_data = df.to_dict(orient='index')
        print(f"INFO: Inventory loaded. {len(existing_data)} records from '{xlsx_filepath}'.")
    except Exception as e:
        print(f"CRITICAL_ERROR loading XLSX from '{xlsx_filepath}': {e}. Inventory will be rebuilt."); traceback.print_exc()
        return {}
    return existing_data

def create_backup(xlsx_filepath, max_backups=5):
    try:
        if os.path.exists(xlsx_filepath):
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_path = f"{xlsx_filepath}.bak.{timestamp}"
            
            # Get existing backups
            existing_backups = [f for f in os.listdir(os.path.dirname(xlsx_filepath)) 
                              if f.startswith(os.path.basename(xlsx_filepath) + ".bak")]
            
            # Remove oldest if too many
            if len(existing_backups) >= max_backups:
                oldest_backup = min(existing_backups, key=lambda x: os.path.getctime(
                    os.path.join(os.path.dirname(xlsx_filepath), x)))
                os.remove(os.path.join(os.path.dirname(xlsx_filepath), oldest_backup))
            
            shutil.copy2(xlsx_filepath, backup_path)
            return True
    except Exception as e:
        print(f"Backup creation failed: {e}")
    return False

def save_inventory_to_xlsx(data_to_save, xlsx_filepath):
    if not isinstance(data_to_save, list):
        print("ERROR: Data to save is not a list.")
        return False
    
    try:
        # Verify data integrity before save
        df = pd.DataFrame(data_to_save if data_to_save else [], columns=FIELDNAMES)
        
        # Critical check: Ensure we're not losing notes
        if os.path.exists(xlsx_filepath):
            try:
                existing_df = pd.read_excel(xlsx_filepath)
                if 'Manual_Notes' in existing_df.columns and 'Full Path' in existing_df.columns:
                    # Create clean mapping of existing notes
                    existing_notes_dict = existing_df.set_index('Full Path')['Manual_Notes'].to_dict()
                    
                    # Set index for current data
                    df = df.set_index('Full Path')
                    
                    # Update notes only where they're empty in current data
                    for idx in df.index:
                        current_note = str(df.at[idx, 'Manual_Notes']).strip()
                        if not current_note and idx in existing_notes_dict:
                            df.at[idx, 'Manual_Notes'] = existing_notes_dict[idx]
                    
                    df = df.reset_index()
            except Exception as e:
                print(f"Warning: Error merging existing notes: {e}")
                if 'Full Path' in df.index.names:
                    df = df.reset_index()

        # Ensure Manual_Notes column exists and is properly formatted
        df['Manual_Notes'] = df['Manual_Notes'].fillna('').astype(str)
        
        # Create backup before saving
        if os.path.exists(xlsx_filepath):
            backup_path = xlsx_filepath + ".bak"
            try:
                import shutil
                shutil.copy2(xlsx_filepath, backup_path)
            except Exception as e:
                print(f"Warning: Backup creation failed: {e}")

        # Use proper temporary file
        temp_filepath = xlsx_filepath + "_temp.xlsx"
        
        # Save to temporary file
        try:
            with pd.ExcelWriter(temp_filepath, engine='openpyxl') as writer:
                df.to_excel(writer, index=False, sheet_name='File Inventory')
                worksheet = writer.sheets['File Inventory']
                for idx, col in enumerate(df.columns, 1):
                    max_length = max(
                        df[col].astype(str).str.len().max(),
                        len(str(col))
                    ) + 2
                    worksheet.column_dimensions[get_column_letter(idx)].width = min(max_length, 70)
            
            # Verify temp file was written correctly
            if os.path.exists(temp_filepath) and os.path.getsize(temp_filepath) > 0:
                # Read back temp file to verify
                pd.read_excel(temp_filepath)
                
                # If verification passed, replace original
                if os.path.exists(xlsx_filepath):
                    os.remove(xlsx_filepath)
                os.rename(temp_filepath, xlsx_filepath)
                print(f"Inventory saved successfully: {len(df)} entries to '{xlsx_filepath}'")
                return True
            
            return False
            
        finally:
            # Clean up temp file if it exists
            if os.path.exists(temp_filepath):
                try:
                    os.remove(temp_filepath)
                except:
                    pass
                    
    except Exception as e:
        print(f"CRITICAL ERROR saving XLSX to '{xlsx_filepath}': {e}")
        traceback.print_exc()
        return False

def process_folder_inventory(start_folder_path, xlsx_output_path, topic_keywords_config):
    if not os.path.isdir(start_folder_path):
        return [], f"Error: Start folder '{start_folder_path}' not found.", 0, 0, 0
    existing_inventory_map = load_existing_inventory(xlsx_output_path)
    paths_from_old_inventory = set(existing_inventory_map.keys())
    current_inventory_list = []
    file_count, updates_count, adds_count, removed_count = 0, 0, 0, 0
    for root, dirs, files in os.walk(start_folder_path, topdown=True):
        dirs[:] = [d for d in dirs if d not in ['.git', '__pycache__', '.ipynb_checkpoints', '.DS_Store']]
        if os.path.abspath(root) == os.path.dirname(os.path.abspath(xlsx_output_path)):
            files = [f for f in files if os.path.basename(f) != os.path.basename(xlsx_output_path)]
        for filename in files:
            if filename.startswith("~$"): continue
            filepath = os.path.join(root, filename)
            abs_filepath = os.path.abspath(filepath)
            try:
                stat_info = os.stat(filepath)
                ext = os.path.splitext(filename)[1].lower()
                current_file_data = {key: ('' if key == 'Manual_Notes' else pd.NA) for key in FIELDNAMES} # Initialize with FIELDNAMES
                current_file_data.update({
                    'Folder Path': root, 'File Name': filename, 'Extension': ext,
                    'Size (Bytes)': stat_info.st_size, 
                    'Last Modified': datetime.datetime.fromtimestamp(stat_info.st_mtime).isoformat(),
                    'Full Path': abs_filepath, 
                    'Content Hint': get_content_hint(filepath, ext),
                    'Identified Topics (DOCX)': "N/A", 'Status': 'Active', 'Manual_Notes': ''
                })
                if ext == '.docx': current_file_data['Identified Topics (DOCX)'] = check_docx_for_topics(filepath, topic_keywords_config)
                if abs_filepath in existing_inventory_map:
                    old_data = existing_inventory_map[abs_filepath]
                    current_file_data['Manual_Notes'] = old_data.get('Manual_Notes', '') 
                    if old_data.get('Size (Bytes)') != stat_info.st_size or old_data.get('Last Modified') != current_file_data['Last Modified']:
                        current_file_data['Status'] = 'Updated'; updates_count += 1
                    paths_from_old_inventory.discard(abs_filepath)
                else:
                    current_file_data['Status'] = 'Added'; adds_count += 1
                current_inventory_list.append(current_file_data)
                file_count += 1
            except Exception as e: print(f"Warning: Could not process file '{filepath}': {e}. Skipping."); pass
    for old_path in paths_from_old_inventory:
        if old_path in existing_inventory_map:
            removed_file_data = {key: val for key, val in existing_inventory_map[old_path].items() if key in FIELDNAMES}
            removed_file_data['Status'] = 'Removed (Not Found)'
            if pd.isna(removed_file_data.get('Manual_Notes')): removed_file_data['Manual_Notes'] = ''
            current_inventory_list.append(removed_file_data)
            removed_count += 1
    status_message = f"Scan Complete. Found {file_count} files. ({adds_count} new, {updates_count} updated, {removed_count} removed)."
    return current_inventory_list, status_message, adds_count, updates_count, removed_count

# --- Gradio Callbacks ---
def shutdown_server():
    try:
        # First inform user
        gr.Info("Shutting down gracefully...")
        time.sleep(0.5)  # Give UI time to show message
        
        # Create backup if needed
        if hasattr(demo, 'current_xlsx_path_state') and demo.current_xlsx_path_state:
            create_backup(demo.current_xlsx_path_state)
            time.sleep(0.5)  # Give time for backup
            
        # Use more graceful shutdown
        def delayed_exit():
            time.sleep(1)  # Wait for response to be sent
            os._exit(0)    # More graceful than os.kill
            
        threading.Thread(target=delayed_exit, daemon=True).start()
        return "Shutdown initiated..."
    except Exception as e:
        print(f"Error during shutdown: {e}")
        os._exit(1)  # Emergency exit if graceful shutdown fails

def open_containing_folder_os(path_to_item):
    if not path_to_item or not isinstance(path_to_item, str): return "Error: Invalid path."
    if not os.path.exists(path_to_item): return f"Error: Path does not exist: {path_to_item}"
    folder_to_open = os.path.dirname(os.path.abspath(path_to_item))
    try:
        if platform.system() == "Windows": os.startfile(folder_to_open)
        elif platform.system() == "Darwin": subprocess.run(["open", folder_to_open], check=True)
        else: subprocess.run(["xdg-open", folder_to_open], check=True)
        return f"Opened: {folder_to_open}"
    except Exception as e: return f"Error opening folder '{folder_to_open}': {e}"

def handle_action_click(current_df_displayed: pd.DataFrame, evt: gr.SelectData):
    try:
        if evt is None or evt.index is None or not hasattr(evt, 'index'):
            return "No cell selected."
        clicked_column_name = current_df_displayed.columns[evt.index[1]]
        if clicked_column_name == 'Action':
            path_to_open = current_df_displayed.iloc[evt.index[0]]['Full Path']
            return open_containing_folder_os(path_to_open)
        return None
    except Exception as e:
        print(f"CRITICAL ERROR in handle_action_click: {e}"); traceback.print_exc()
        return f"Error handling action: {e}"

def run_scan_and_display(folder_path_input):
    try:
        empty_df_for_display = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
        empty_full_df_for_state = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action':'📂', 'Manual_Notes':''})

        if not folder_path_input or not os.path.isdir(folder_path_input):
            return empty_df_for_display, "Error: Invalid folder path.", "", empty_full_df_for_state, ""
        
        xlsx_for_this_folder = os.path.join(folder_path_input, INVENTORY_FILENAME)
        print(f"Starting scan: {folder_path_input}. Inventory: {xlsx_for_this_folder}")
        
        # Add recovery attempt
        if os.path.exists(xlsx_for_this_folder):
            attempt_inventory_recovery(xlsx_for_this_folder)
        
        inventory_data_list, status_msg, _, _, _ = process_folder_inventory(folder_path_input, xlsx_for_this_folder, TOPIC_KEYWORDS)
        
        df_full_state_candidate = empty_full_df_for_state.copy()
        if inventory_data_list: 
            df_temp = pd.DataFrame(inventory_data_list) # Contains FIELDNAMES columns
            # Merge with an empty df that has all STATE_COLUMNS to ensure structure
            df_full_state_candidate = pd.concat([empty_full_df_for_state, df_temp], ignore_index=True)
            df_full_state_candidate = df_full_state_candidate.reindex(columns=STATE_COLUMNS) # Ensure order
            df_full_state_candidate['Action'] = '📂' # Add/ensure icon
            df_full_state_candidate['Manual_Notes'] = df_full_state_candidate['Manual_Notes'].fillna('')
            
            if not df_full_state_candidate.empty:
                # Save only data columns (FIELDNAMES) to Excel
                save_inventory_to_xlsx(df_full_state_candidate[FIELDNAMES].to_dict(orient='records'), xlsx_for_this_folder)
        
        # Prepare display_df from df_full_state_candidate
        display_df = df_full_state_candidate.reindex(columns=DISPLAY_COLUMNS).fillna('')
        
        return display_df, status_msg, xlsx_for_this_folder, df_full_state_candidate.copy(), xlsx_for_this_folder
    except Exception as e:
        print(f"CRITICAL ERROR in run_scan_and_display: {e}"); traceback.print_exc()
        _empty_display = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
        _empty_state = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action':'📂', 'Manual_Notes':''})
        return _empty_display, f"Error during scan: {e}", "", _empty_state, ""

def filter_dataframe_display(df_full_from_state: pd.DataFrame, status_filter: str, topic_filter_text: str, filename_filter_text: str) -> pd.DataFrame:
    try:
        empty_df_with_display_columns = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
        
        if not isinstance(df_full_from_state, pd.DataFrame) or df_full_from_state.empty:
            return empty_df_with_display_columns

        df_filtered = df_full_from_state.copy()

        # Status filter (unchanged)
        if status_filter and status_filter != "All":
            if 'Status' in df_filtered.columns:
                df_filtered = df_filtered[df_filtered['Status'] == status_filter]
            else:
                return empty_df_with_display_columns.copy()
        
        # Handle filename filter with folder inclusions/exclusions
        if filename_filter_text:
            if 'File Name' not in df_filtered.columns or 'Folder Path' not in df_filtered.columns:
                return empty_df_with_display_columns.copy()
                
            # Split into folder filters and filename filters
            folder_excludes = []
            folder_includes = []
            filename_terms = []
            
            for term in filename_filter_text.split(','):
                term = term.strip().lower()
                if term.startswith('folder:'):
                    # Folder exclusion
                    folder_term = term[7:].strip()
                    if folder_term:
                        folder_excludes.append(folder_term)
                elif term.startswith('incfolder:'):
                    # Folder inclusion
                    folder_term = term[10:].strip()
                    if folder_term:
                        folder_includes.append(folder_term)
                else:
                    if term:
                        filename_terms.append(term)
            
            # Apply folder inclusions (OR logic - keep if matches any term)
            if folder_includes:
                mask = pd.Series(False, index=df_filtered.index)
                for term in folder_includes:
                    mask |= df_filtered['Folder Path'].astype(str).str.lower().str.contains(term, na=False)
                df_filtered = df_filtered[mask]
            
            # Apply folder exclusions (OR logic - exclude if matches any term)
            if folder_excludes:
                for term in folder_excludes:
                    df_filtered = df_filtered[
                        ~df_filtered['Folder Path'].astype(str).str.lower().str.contains(term, na=False)
                    ]
            
            # Apply filename filters (AND logic)
            if filename_terms:
                for term in filename_terms:
                    df_filtered = df_filtered[
                        df_filtered['File Name'].astype(str).str.lower().str.contains(term, na=False)
                    ]
                
        # Topic filter (unchanged)
        if topic_filter_text:
            if 'Identified Topics (DOCX)' in df_filtered.columns:
                search_terms = [term.strip().lower() for term in topic_filter_text.split(',')]
                for term in search_terms:
                    df_filtered = df_filtered[
                        df_filtered['Identified Topics (DOCX)'].astype(str).str.lower().str.contains(term, na=False)
                    ]
                
        return df_filtered.reindex(columns=DISPLAY_COLUMNS).fillna('')
    except Exception as e:
        print(f"CRITICAL ERROR in filter_dataframe_display: {e}")
        traceback.print_exc()
        return pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')


def save_notes(displayed_df_with_edits: pd.DataFrame, full_df_from_state: pd.DataFrame, xlsx_path: str):
    try:
        empty_state_df = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action':'📂', 'Manual_Notes':''})
        if not isinstance(full_df_from_state, pd.DataFrame) or full_df_from_state.empty:
            return full_df_from_state.copy() if isinstance(full_df_from_state, pd.DataFrame) else empty_state_df, "Cannot save: Master data is empty."
        if not isinstance(displayed_df_with_edits, pd.DataFrame):
            return full_df_from_state.copy(), "No data displayed to save from."

        updated_full_df = full_df_from_state.copy() # Has 'Action' column
        if 'Full Path' not in updated_full_df.columns:
             print("CRITICAL: 'Full Path' column missing in full_df_from_state for save_notes.")
             return full_df_from_state.copy(), "Error: 'Full Path' column missing."
        
        updated_full_df.set_index('Full Path', inplace=True, drop=False) # Keep 'Full Path' as column too for reindex
        
        if not displayed_df_with_edits.empty and \
           'Full Path' in displayed_df_with_edits.columns and \
           'Manual_Notes' in displayed_df_with_edits.columns:
            notes_to_update_df = displayed_df_with_edits.set_index('Full Path')[['Manual_Notes']].fillna('')
            updated_full_df.update(notes_to_update_df) # Updates 'Manual_Notes' based on index
        
        updated_full_df.reset_index(drop=True, inplace=True) # Remove old index, keep 'Full Path' column
        
        # Ensure final state df has all STATE_COLUMNS and 'Action' is filled
        final_df_for_state = updated_full_df.reindex(columns=STATE_COLUMNS)
        final_df_for_state['Action'] = final_df_for_state['Action'].fillna('📂')
        final_df_for_state['Manual_Notes'] = final_df_for_state['Manual_Notes'].fillna('')

        status_msg = "No valid path."
        if xlsx_path and isinstance(xlsx_path, str) and xlsx_path.strip():
            # Save only FIELDNAMES (data columns) to Excel
            if save_inventory_to_xlsx(final_df_for_state[FIELDNAMES].to_dict(orient='records'), xlsx_path):
                timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                status_msg = f"Notes saved successfully to {os.path.basename(xlsx_path)}. (at {timestamp})"
            else: status_msg = "Error: Failed to save notes to file."
        else: status_msg = "Error: Inventory file path not set."
        
        return final_df_for_state, status_msg, f"Last saved: {timestamp}"
    except Exception as e:
        print(f"CRITICAL ERROR in save_notes: {e}"); traceback.print_exc()
        _empty_state = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action':'📂', 'Manual_Notes':''})
        return full_df_from_state.copy() if isinstance(full_df_from_state, pd.DataFrame) else _empty_state, f"Error saving notes: {e}", "Last saved: Error"

def attempt_inventory_recovery(xlsx_filepath):
    try:
        temp_path = xlsx_filepath + ".xlsx.tmp"
        backup_path = xlsx_filepath + ".bak"
        
        # Only recover if the main file is missing or empty
        if not os.path.exists(xlsx_filepath) or os.path.getsize(xlsx_filepath) == 0:
            # First try temp file
            if os.path.exists(temp_path) and os.path.getsize(temp_path) > 0:
                try:
                    temp_df = pd.read_excel(temp_path)
                    if not temp_df.empty:
                        os.replace(temp_path, xlsx_filepath)
                        print("Recovered from temporary save file")
                        return True
                except:
                    if os.path.exists(temp_path):
                        os.remove(temp_path)
            
            # Then try backup
            if os.path.exists(backup_path) and os.path.getsize(backup_path) > 0:
                try:
                    backup_df = pd.read_excel(backup_path)
                    if not backup_df.empty:
                        import shutil
                        shutil.copy2(backup_path, xlsx_filepath)
                        print("Recovered from backup file")
                        return True
                except:
                    pass
        return False
    except Exception as e:
        print(f"Recovery attempt failed: {e}")
        return False

class ErrorTracker:
    def __init__(self):
        self.error_count = 0
        self.last_error_time = None
        self.max_errors = 3
        self.error_window = 300  # 5 minutes

    def record_error(self):
        now = time.time()
        if self.last_error_time and (now - self.last_error_time) > self.error_window:
            self.error_count = 0
        self.error_count += 1
        self.last_error_time = now
        return self.error_count >= self.max_errors

error_tracker = ErrorTracker()

# --- Gradio UI ---
with gr.Blocks(theme=gr.themes.Soft()) as demo:
    # State holds data + 'Action' column for internal consistency
    # Initial empty state must match this structure
    _initial_empty_state_df = pd.DataFrame(columns=STATE_COLUMNS)
    _initial_empty_state_df['Action'] = '📂' # Ensure Action column has icon even if 0 rows
    _initial_empty_state_df['Manual_Notes'] = ''


    full_df_state = gr.State(_initial_empty_state_df)
    current_xlsx_path_state = gr.State("")

    gr.Markdown(f"# File Inventory Dashboard v{VERSION}")
    gr.Markdown("""<p style='font-size: 0.9em; color: #666; margin-bottom: 10px;'>
        <em>Quick Guide:</em> 📂=open folder | Edit 'Manual_Notes' directly | Search: terms with comma | 
        Filters: folder:old (exclude), IncFolder:data (include)</p>""")
    with gr.Row():
        folder_input = gr.Textbox(label="Folder Path", placeholder="Enter folder path", value=START_FOLDER)
        scan_button = gr.Button("Scan Folder & Load/Update Inventory", variant="primary")
        shutdown_button = gr.Button("Exit Application", variant="stop")

    status_output = gr.Markdown()
    file_op_status_output = gr.Markdown()
    current_xlsx_display = gr.Textbox(label="Active Inventory File Path", interactive=False)
    
    gr.Markdown("## Filter and Save")
    with gr.Row():
        status_dropdown = gr.Dropdown(choices=["All", "Active", "Updated", "Added", "Removed (Not Found)"], value="All", label="Filter by Status")
        topic_search = gr.Textbox(label="Search Topics")
        filename_search = gr.Textbox(
            label="Search File Name", 
            placeholder="Examples: manuscript, draft | folder:old (exclude) | IncFolder:data (include only)"
        )
    with gr.Row():
        filter_button = gr.Button("Apply Filters")
        save_notes_button = gr.Button("Save Notes", variant="primary")

    gr.Markdown("## Inventory Data (Click '📂' to open folder, edit 'Manual_Notes' column, then click 'Save Notes')")
    
    # dataframe_output displays DISPLAY_COLUMNS
    _initial_empty_display_df = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
    dataframe_output = gr.DataFrame(
        value=_initial_empty_display_df, 
        interactive=True, wrap=True, headers=DISPLAY_COLUMNS,
        datatype=['markdown'] + ['str']*(len(DISPLAY_COLUMNS)-1)
    )

    with gr.Row():
        last_save_indicator = gr.Markdown("Last saved: Never")

    scan_button.click(
        fn=run_scan_and_display,
        inputs=[folder_input],
        outputs=[dataframe_output, status_output, current_xlsx_display, full_df_state, current_xlsx_path_state]
    )
    filter_button.click(
        fn=filter_dataframe_display,
        inputs=[full_df_state, status_dropdown, topic_search, filename_search],
        outputs=[dataframe_output]
    )
    shutdown_button.click(
        fn=shutdown_server,
        inputs=None,
        outputs=gr.Textbox(visible=False)
    )
    dataframe_output.select(
        fn=handle_action_click, inputs=[dataframe_output], outputs=[file_op_status_output]
    )
    save_notes_button.click(
        fn=save_notes,
        inputs=[dataframe_output, full_df_state, current_xlsx_path_state],
        outputs=[full_df_state, file_op_status_output, last_save_indicator]
    )

def setup_shutdown_handler():
    import atexit
    
    def cleanup():
        print("Performing cleanup before exit...")
        # Save any pending changes
        if hasattr(demo, 'current_xlsx_path_state') and demo.current_xlsx_path_state:
            create_backup(demo.current_xlsx_path_state)
    
    atexit.register(cleanup)

# Add to main:
if __name__ == "__main__":
    try:
        setup_shutdown_handler()
        logging.info(f"Launching File Inventory Dashboard v{VERSION}")
        logging.info(f"Start folder: {START_FOLDER}")
        demo.launch(
            server_name="127.0.0.1", 
            server_port=7860, 
            inbrowser=True,
            show_error=True
        )
    except Exception as e:
        logging.error(f"Critical startup error: {e}")
        traceback.print_exc()