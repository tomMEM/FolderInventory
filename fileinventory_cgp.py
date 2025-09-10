"""
File Inventory Dashboard
Version: 6.2
Author: TB (updated collaboratively)
Description: A Gradio-based dashboard for managing file inventories with search, filtering, and note-taking capabilities.
Last Updated: 2025-09-10

Key improvements in 6.1 (no new libraries, structure preserved):
- Stronger Manual_Notes stability: notes are preserved/merged on every save.
- Removed/deleted files are retained in the inventory if they have Manual_Notes.
- Searching now matches across key columns, including Manual_Notes (and others) without changing the UI.
- Safer backups, recovery, and temporary save/replace logic.
- Minor robustness fixes (imports, timestamp handling, edge cases).

New in 6.2 (small, focused change):
- Folder Path input is now a Dropdown with memory (type new or choose from history).
"""

# ==============================================================================
# FINAL, STABLE, AND WORKING SCRIPT (VERSION 6.2)
# ==============================================================================

import os
import sys
import time
import threading
import datetime
import openpyxl
from openpyxl.utils import get_column_letter
import platform
import subprocess
import gradio as gr
import pandas as pd
import webbrowser
import traceback
import logging
import shutil  # used for backups/copies

# --- PyInstaller Windowed Mode Fix ---
if sys.stderr is None:
    class DummyStream:
        def isatty(self): return False
        def write(self, msg): pass
        def flush(self): pass
    sys.stderr = DummyStream()
    sys.stdout = DummyStream()

# --- Proxy Setup (omitted for brevity, assume it's correct) ---
PROXY_URL = os.environ.get("HTTP_PROXY_URL", "http://localhost:4321")
if PROXY_URL:
    os.environ['http_proxy'] = PROXY_URL
    os.environ['https_proxy'] = PROXY_URL
else:
    pass

current_no_proxy = os.environ.get('NO_PROXY', '')
additional_no_proxy_hosts = ['localhost', '127.0.0.1', '0.0.0.0']
new_no_proxy_parts = [host for host in current_no_proxy.split(',') if host.strip()]
for host in additional_no_proxy_hosts:
    if host not in new_no_proxy_parts:
        new_no_proxy_parts.append(host)
os.environ['NO_PROXY'] = ','.join(new_no_proxy_parts)

# --- Configuration ---
VERSION = "6.2.0"
BUILD_DATE = "2025-09-10"
START_FOLDER = r"\\synology\YanYan\TB\1.Manuscript"
INVENTORY_FILENAME = "inventory.xlsx"

# Recent folders memory (no new libs; plain text file)
RECENT_FOLDERS_FILE = os.path.join(os.path.dirname(os.path.abspath(__file__)), "recent_folders.txt")
MAX_RECENT_FOLDERS = 15

# Topic keywords configuration
TOPIC_KEYWORDS = {
    "PET": {
        "all_required": ["pet"],
        "any_of_these": ["scan", "imaging", "tracer"]
    },
    "DID": {
        "all_required": ["did"],
        "any_of_these": ["dementia", "cognitive", "decline"]
    },
    "AD": {
        "all_required": ["alzheimer"],
        "any_of_these": ["disease", "dementia", "ad"]
    }
}

# File type configurations
TEXT_BASED_EXTENSIONS = [
    '.docx', '.pptx', '.txt',
    '.py', '.r', '.md'
]
SPREADSHEET_EXTENSIONS = [
    '.xlsx', '.csv', '.prism'
]

# Column definitions
FIELDNAMES = [
    'Folder Path', 'File Name', 'Extension',
    'Size (Bytes)', 'Last Modified', 'Full Path',
    'Content Hint', 'Identified Topics (DOCX)',
    'Status', 'Manual_Notes'
]

DISPLAY_COLUMNS = [
    'Action', 'File Name', 'Status',
    'Last Modified', 'Identified Topics (DOCX)',
    'Content Hint', 'Manual_Notes', 'Full Path'
]

STATE_COLUMNS = FIELDNAMES + ['Action']  # Columns for the full_df_state

# Error handling configuration
MAX_BACKUP_FILES = 5
ERROR_WINDOW_SECONDS = 300
MAX_ERROR_COUNT = 3

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('inventory_dashboard.log'),
        logging.StreamHandler()
    ]
)


# =========================
# Helper Functions
# =========================

def get_content_hint(filepath, extension):
    """
    Produces a lightweight content hint without new dependencies.
    """
    hint = "N/A"
    try:
        if extension == '.docx':
            try:
                import docx
                doc = docx.Document(filepath)
                if doc.paragraphs:
                    hint = "First para: " + doc.paragraphs[0].text[:150] + "..."
                else:
                    hint = "DOCX: No paragraphs found."
            except ImportError:
                hint = "DOCX: 'python-docx' not installed."
            except Exception:
                hint = "DOCX: Corrupt or unreadable."
        elif extension == '.pptx':
            try:
                import pptx
                prs = pptx.Presentation(filepath)
                if prs.slides and prs.slides[0].shapes.title:
                    hint = "First slide title: " + prs.slides[0].shapes.title.text[:150]
                elif prs.slides:
                    hint = "PPTX: First slide no title."
                else:
                    hint = "PPTX: No slides."
            except ImportError:
                hint = "PPTX: 'python-pptx' not installed."
            except Exception:
                hint = "PPTX: Corrupt or unreadable."
        elif extension in TEXT_BASED_EXTENSIONS:
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    lines = [line.strip() for i, line in enumerate(f) if i < 2]
                    hint_text = " ".join(lines)
                    hint = ("First 2 lines: " + hint_text[:200] + "...") if hint_text else f"{extension.upper()}: Empty"
            except Exception as e:
                hint = f"{extension.upper()}: Read error"
        elif extension in SPREADSHEET_EXTENSIONS:
            hint = "Spreadsheet file."
    except Exception as e:
        hint = f"Hint Error for {os.path.basename(filepath)}: {e}"
    return hint


def check_docx_for_topics(filepath, topic_definitions):
    identified_topics = []
    try:
        # Check if file exists and is accessible first
        if not os.path.exists(filepath):
            print(f"Warning: DOCX file not found: {filepath}")
            return "N/A (File not found)"

        import docx
        try:
            doc = docx.Document(filepath)
        except Exception as e:
            print(f"Warning: Could not open DOCX {filepath}: {e}")
            return "N/A (Access error)"

        full_text = "\n".join([para.text for para in doc.paragraphs]).lower()
        if not full_text.strip():
            return "DOCX Empty"

        for topic_name, criteria in topic_definitions.items():
            all_match = all(kw.lower() in full_text for kw in criteria.get("all_required", []))
            any_match = True
            if "any_of_these" in criteria and criteria["any_of_these"]:
                any_match = any(kw.lower() in full_text for kw in criteria["any_of_these"])
            if all_match and any_match:
                identified_topics.append(topic_name)
    except Exception as e:
        print(f"Error checking DOCX topics for {filepath}: {e}")
        return "N/A (Error reading DOCX)"
    return ", ".join(identified_topics) if identified_topics else "N/A"


def _ensure_expected_columns(df: pd.DataFrame) -> pd.DataFrame:
    df = df.copy()
    for col_name in FIELDNAMES:
        if col_name not in df.columns:
            df[col_name] = '' if col_name == 'Manual_Notes' else pd.NA
    if 'Manual_Notes' in df.columns:
        df['Manual_Notes'] = df['Manual_Notes'].fillna('').astype(str)
    if 'Size (Bytes)' in df.columns:
        df['Size (Bytes)'] = pd.to_numeric(df['Size (Bytes)'], errors='coerce').astype('Int64')
    return df



def load_existing_inventory(xlsx_filepath):
    """
    Load existing inventory into a dict keyed by 'Full Path'.
    Missing columns are added. Manual_Notes coerced to string.
    """
    existing_data = {}
    if not os.path.exists(xlsx_filepath):
        print(f"INFO: Inventory file '{xlsx_filepath}' not found. New one will be created.")
        return existing_data
    print(f"INFO: Attempting to load inventory from '{xlsx_filepath}'...")
    try:
        df = pd.read_excel(xlsx_filepath, engine='openpyxl')
        if 'Full Path' not in df.columns:
            print(f"ERROR: 'Full Path' column missing in '{xlsx_filepath}'. Cannot process.")
            return {}
        df = _ensure_expected_columns(df)
        df.set_index('Full Path', inplace=True)
        existing_data = df.to_dict(orient='index')
        print(f"INFO: Inventory loaded. {len(existing_data)} records from '{xlsx_filepath}'.")
    except Exception as e:
        print(f"CRITICAL_ERROR loading XLSX from '{xlsx_filepath}': {e}. Inventory will be rebuilt.")
        traceback.print_exc()
        return {}
    return existing_data


def create_backup(xlsx_filepath, max_backups=5):
    """
    Create rotating backups inventory.xlsx.bak.YYYYMMDD_HHMMSS (up to max_backups).
    """
    try:
        if os.path.exists(xlsx_filepath):
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_path = f"{xlsx_filepath}.bak.{timestamp}"

            # Get existing backups
            dir_name = os.path.dirname(xlsx_filepath) or "."
            prefix = os.path.basename(xlsx_filepath) + ".bak"
            existing_backups = [f for f in os.listdir(dir_name) if f.startswith(prefix)]

            # Remove oldest if too many
            if len(existing_backups) >= max_backups:
                oldest_backup = min(existing_backups, key=lambda x: os.path.getctime(
                    os.path.join(dir_name, x)))
                os.remove(os.path.join(dir_name, oldest_backup))

            shutil.copy2(xlsx_filepath, backup_path)
            return True
    except Exception as e:
        print(f"Backup creation failed: {e}")
    return False


def _merge_existing_notes(df_new: pd.DataFrame, xlsx_filepath: str) -> pd.DataFrame:
    """
    Merge Manual_Notes from existing file into df_new based on 'Full Path'.
    Does not overwrite non-empty notes in df_new.
    """
    df = df_new.copy()
    try:
        if os.path.exists(xlsx_filepath):
            existing_df = pd.read_excel(xlsx_filepath, engine='openpyxl')
            if 'Full Path' in existing_df.columns and 'Manual_Notes' in existing_df.columns:
                existing_df = _ensure_expected_columns(existing_df)
                existing_notes = existing_df.set_index('Full Path')['Manual_Notes'].to_dict()

                df = df.set_index('Full Path')
                if 'Manual_Notes' not in df.columns:
                    df['Manual_Notes'] = ''
                df['Manual_Notes'] = df['Manual_Notes'].fillna('').astype(str)

                for idx in df.index:
                    current_note = str(df.at[idx, 'Manual_Notes']).strip()
                    if not current_note and idx in existing_notes:
                        df.at[idx, 'Manual_Notes'] = existing_notes[idx]

                df = df.reset_index()
    except Exception as e:
        print(f"Warning: Error merging existing notes: {e}")
    return df


def save_inventory_to_xlsx(data_to_save, xlsx_filepath):
    """
    Save inventory to XLSX with:
    - Merge of existing Manual_Notes (never lose notes).
    - Temp write + verify + replace.
    - Column width autosizing.
    """
    if not isinstance(data_to_save, list):
        print("ERROR: Data to save is not a list.")
        return False

    try:
        df = pd.DataFrame(data_to_save if data_to_save else [], columns=FIELDNAMES)
        df = _ensure_expected_columns(df)

        # Merge existing notes so we never lose them
        df = _merge_existing_notes(df, xlsx_filepath)

        # Ensure Manual_Notes column exists and is properly formatted
        df['Manual_Notes'] = df['Manual_Notes'].fillna('').astype(str)

        # Create simple rolling backup (non-rotating) in addition to rotating backups
        if os.path.exists(xlsx_filepath):
            backup_path = xlsx_filepath + ".bak"
            try:
                shutil.copy2(xlsx_filepath, backup_path)
            except Exception as e:
                print(f"Warning: Backup creation failed: {e}")

        # Use proper temporary file
        temp_filepath = xlsx_filepath + "_temp.xlsx"

        # Save to temporary file
        try:
            with pd.ExcelWriter(temp_filepath, engine='openpyxl') as writer:
                df.to_excel(writer, index=False, sheet_name='File Inventory')
                worksheet = writer.sheets['File Inventory']
                for idx, col in enumerate(df.columns, 1):
                    max_length = max(
                        df[col].astype(str).str.len().max(),
                        len(str(col))
                    )
                    if pd.isna(max_length):
                        max_length = len(str(col))
                    worksheet.column_dimensions[get_column_letter(idx)].width = min(int(max_length) + 2, 70)

            # Verify temp file was written correctly
            if os.path.exists(temp_filepath) and os.path.getsize(temp_filepath) > 0:
                # Read back temp file to verify
                pd.read_excel(temp_filepath, engine='openpyxl')

                # If verification passed, replace original
                if os.path.exists(xlsx_filepath):
                    os.remove(xlsx_filepath)
                os.rename(temp_filepath, xlsx_filepath)
                print(f"Inventory saved successfully: {len(df)} entries to '{xlsx_filepath}'")
                return True

            return False

        finally:
            # Clean up temp file if it exists (if we didn't rename it)
            if os.path.exists(temp_filepath):
                try:
                    os.remove(temp_filepath)
                except:
                    pass

    except Exception as e:
        print(f"CRITICAL ERROR saving XLSX to '{xlsx_filepath}': {e}")
        traceback.print_exc()
        return False


def process_folder_inventory(start_folder_path, xlsx_output_path, topic_keywords_config):
    """
    Build the current file list, detect Added/Updated, and optionally carry forward
    Removed items that have Manual_Notes (to preserve user-entered notes).
    """
    if not os.path.isdir(start_folder_path):
        return [], f"Error: Start folder '{start_folder_path}' not found.", 0, 0, 0

    existing_inventory_map = load_existing_inventory(xlsx_output_path)
    paths_from_old_inventory = set(existing_inventory_map.keys())

    current_inventory_list = []
    file_count, updates_count, adds_count, removed_count = 0, 0, 0, 0

    for root, dirs, files in os.walk(start_folder_path, topdown=True):
        files = [f for f in files if f.lower() != INVENTORY_FILENAME.lower()]

        # Avoid listing the output inventory file itself
        if os.path.abspath(root) == os.path.dirname(os.path.abspath(xlsx_output_path)):
            files = [f for f in files if os.path.basename(f) != os.path.basename(xlsx_output_path)]

        for filename in files:
            if filename.startswith("~$"):  # Skip Office temp files
                continue

            filepath = os.path.join(root, filename)
            abs_filepath = os.path.abspath(filepath)
            try:
                stat_info = os.stat(filepath)
                ext = os.path.splitext(filename)[1].lower()

                # Initialize with FIELDNAMES
                current_file_data = {key: ('' if key == 'Manual_Notes' else pd.NA) for key in FIELDNAMES}
                current_file_data.update({
                    'Folder Path': root,
                    'File Name': filename,
                    'Extension': ext,
                    'Size (Bytes)': stat_info.st_size,
                    'Last Modified': datetime.datetime.fromtimestamp(stat_info.st_mtime).isoformat(),
                    'Full Path': abs_filepath,
                    'Content Hint': get_content_hint(filepath, ext),
                    'Identified Topics (DOCX)': "N/A",
                    'Status': 'Active',
                    'Manual_Notes': ''
                })

                if ext == '.docx':
                    current_file_data['Identified Topics (DOCX)'] = check_docx_for_topics(filepath, topic_keywords_config)

                if abs_filepath in existing_inventory_map:
                    old_data = existing_inventory_map[abs_filepath]
                    # Always carry forward existing notes
                    current_file_data['Manual_Notes'] = str(old_data.get('Manual_Notes', '') or '')

                    # Detect update
                    old_size = old_data.get('Size (Bytes)')
                    old_mtime = old_data.get('Last Modified')
                    if (pd.notna(old_size) and old_size != stat_info.st_size) or (pd.notna(old_mtime) and old_mtime != current_file_data['Last Modified']):
                        current_file_data['Status'] = 'Updated'
                        updates_count += 1

                    paths_from_old_inventory.discard(abs_filepath)
                else:
                    current_file_data['Status'] = 'Added'
                    adds_count += 1

                current_inventory_list.append(current_file_data)
                file_count += 1

            except Exception as e:
                print(f"Warning: Could not process file '{filepath}': {e}. Skipping.")
                pass

    # Handle removed files: retain ONLY those with Manual_Notes to preserve user writing
    for old_path in list(paths_from_old_inventory):
        old_row = existing_inventory_map.get(old_path, {})
        # Normalize notes
        old_note = str(old_row.get('Manual_Notes', '') or '')
        if old_note.strip():
            # Build a row using existing data (keep prior metadata if present)
            removed_file_data = {key: old_row.get(key, ('' if key == 'Manual_Notes' else pd.NA)) for key in FIELDNAMES}
            removed_file_data['Status'] = 'Removed (Not Found)'
            removed_file_data['Manual_Notes'] = old_note
            # Ensure Full Path is set (index was 'Full Path')
            removed_file_data['Full Path'] = old_path
            current_inventory_list.append(removed_file_data)
            removed_count += 1
        # else: no notes -> drop from inventory to reduce noise

    status_message = f"Scan Complete. Found {file_count} files. ({adds_count} new, {updates_count} updated, {removed_count} removed-kept-with-notes)."
    return current_inventory_list, status_message, adds_count, updates_count, removed_count


# =========================
# Recent Folders Helpers (for the dropdown)
# =========================

def _read_recent_folders():
    try:
        if os.path.exists(RECENT_FOLDERS_FILE):
            with open(RECENT_FOLDERS_FILE, 'r', encoding='utf-8', errors='ignore') as f:
                lines = [ln.strip() for ln in f.readlines() if ln.strip()]
                # Deduplicate preserving order
                seen = set()
                uniq = []
                for p in lines:
                    if p not in seen:
                        seen.add(p)
                        uniq.append(p)
                return uniq[:MAX_RECENT_FOLDERS]
    except Exception as e:
        print(f"Warning reading recent folders: {e}")
    return []

def _write_recent_folders(paths_list):
    try:
        with open(RECENT_FOLDERS_FILE, 'w', encoding='utf-8') as f:
            for p in paths_list[:MAX_RECENT_FOLDERS]:
                f.write(p + '\n')
    except Exception as e:
        print(f"Warning writing recent folders: {e}")

def load_recent_folders(initial_choice=None):
    """
    Returns (choices, value) for the dropdown.
    Ensures START_FOLDER is present at least on first run.
    """
    choices = _read_recent_folders()
    if not choices:
        choices = [START_FOLDER]
    else:
        # Prefer initial_choice or START_FOLDER as selected value if present
        if initial_choice and initial_choice in choices:
            return choices, initial_choice
        if START_FOLDER in choices:
            return choices, START_FOLDER
    return choices, choices[0] if choices else START_FOLDER

def add_recent_folder(path_str):
    """
    Add a folder to the recent list (dedup, move-to-front).
    Only add if it looks like a valid folder path on this machine.
    """
    try:
        p = str(path_str or '').strip()
        if not p:
            return
        if not os.path.isdir(p):
            # do not add invalid paths
            return
        current = _read_recent_folders()
        # move-to-front dedup
        current = [x for x in current if x != p]
        current.insert(0, p)
        _write_recent_folders(current)
    except Exception as e:
        print(f"Warning adding recent folder: {e}")


# =========================
# Gradio Callbacks
# =========================

def shutdown_server():
    try:
        # First inform user
        gr.Info("Shutting down gracefully...")
        time.sleep(0.5)  # Give UI time to show message

        # Create backup if needed
        if hasattr(demo, 'current_xlsx_path_state') and demo.current_xlsx_path_state:
            create_backup(demo.current_xlsx_path_state)
            time.sleep(0.5)  # Give time for backup

        # Use more graceful shutdown
        def delayed_exit():
            time.sleep(1)  # Wait for response to be sent
            os._exit(0)    # More graceful than os.kill

        threading.Thread(target=delayed_exit, daemon=True).start()
        return "Shutdown initiated..."
    except Exception as e:
        print(f"Error during shutdown: {e}")
        os._exit(1)  # Emergency exit if graceful shutdown fails


def open_containing_folder_os(path_to_item):
    if not path_to_item or not isinstance(path_to_item, str):
        return "Error: Invalid path."
    if not os.path.exists(path_to_item):
        return f"Error: Path does not exist: {path_to_item}"
    folder_to_open = os.path.dirname(os.path.abspath(path_to_item))
    try:
        if platform.system() == "Windows":
            os.startfile(folder_to_open)
        elif platform.system() == "Darwin":
            subprocess.run(["open", folder_to_open], check=True)
        else:
            subprocess.run(["xdg-open", folder_to_open], check=True)
        return f"Opened: {folder_to_open}"
    except Exception as e:
        return f"Error opening folder '{folder_to_open}': {e}"


def handle_action_click(current_df_displayed: pd.DataFrame, evt: gr.SelectData):
    try:
        if evt is None or evt.index is None or not hasattr(evt, 'index'):
            return "No cell selected."
        clicked_column_name = current_df_displayed.columns[evt.index[1]]
        if clicked_column_name == 'Action':
            path_to_open = current_df_displayed.iloc[evt.index[0]]['Full Path']
            return open_containing_folder_os(path_to_open)
        return None
    except Exception as e:
        print(f"CRITICAL ERROR in handle_action_click: {e}")
        traceback.print_exc()
        return f"Error handling action: {e}"


def attempt_inventory_recovery(xlsx_filepath):
    try:
        # Recognize both temp naming styles
        temp_path_candidates = [
            xlsx_filepath + "_temp.xlsx",
            xlsx_filepath + ".xlsx.tmp"
        ]
        backup_path = xlsx_filepath + ".bak"

        # Only recover if the main file is missing or empty
        if (not os.path.exists(xlsx_filepath)) or os.path.getsize(xlsx_filepath) == 0:
            # Try temp candidates
            for temp_path in temp_path_candidates:
                if os.path.exists(temp_path) and os.path.getsize(temp_path) > 0:
                    try:
                        temp_df = pd.read_excel(temp_path, engine='openpyxl')
                        if not temp_df.empty:
                            os.replace(temp_path, xlsx_filepath)
                            print("Recovered from temporary save file")
                            return True
                    except:
                        try:
                            os.remove(temp_path)
                        except:
                            pass

            
            # Then try backup
            if os.path.exists(backup_path) and os.path.getsize(backup_path) > 0:
                try:
                    backup_df = pd.read_excel(backup_path, engine='openpyxl')
                    if not backup_df.empty:
                        shutil.copy2(backup_path, xlsx_filepath)
                        print("Recovered from backup file")
                        return True
                except:
                    pass

        return False
    except Exception as e:
        print(f"Recovery attempt failed: {e}")
        return False


def run_scan_and_display(folder_path_input):
    try:
        empty_df_for_display = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
        empty_full_df_for_state = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action': '📂', 'Manual_Notes': ''})

        if not folder_path_input or not os.path.isdir(folder_path_input):
            return empty_df_for_display, "Error: Invalid folder path.", "", empty_full_df_for_state, ""

        xlsx_for_this_folder = os.path.join(folder_path_input, INVENTORY_FILENAME)
        print(f"Starting scan: {folder_path_input}. Inventory: {xlsx_for_this_folder}")

        # Attempt recovery if needed
        if os.path.exists(xlsx_for_this_folder):
            attempt_inventory_recovery(xlsx_for_this_folder)

        inventory_data_list, status_msg, _, _, _ = process_folder_inventory(folder_path_input, xlsx_for_this_folder, TOPIC_KEYWORDS)

        df_full_state_candidate = empty_full_df_for_state.copy()
        if inventory_data_list:
            df_temp = pd.DataFrame(inventory_data_list)  # Contains FIELDNAMES columns
            # Merge with an empty df that has all STATE_COLUMNS to ensure structure
            df_full_state_candidate = pd.concat([empty_full_df_for_state, df_temp], ignore_index=True)
            df_full_state_candidate = df_full_state_candidate.reindex(columns=STATE_COLUMNS)  # Ensure order
            df_full_state_candidate['Action'] = '📂'  # Add/ensure icon
            df_full_state_candidate['Manual_Notes'] = df_full_state_candidate['Manual_Notes'].fillna('')

            if not df_full_state_candidate.empty:
                # Save only data columns (FIELDNAMES) to Excel with merge protection for notes
                save_inventory_to_xlsx(df_full_state_candidate[FIELDNAMES].to_dict(orient='records'), xlsx_for_this_folder)

        # Prepare display_df from df_full_state_candidate
        display_df = df_full_state_candidate.reindex(columns=DISPLAY_COLUMNS).fillna('')

        return display_df, status_msg, xlsx_for_this_folder, df_full_state_candidate.copy(), xlsx_for_this_folder
    except Exception as e:
        print(f"CRITICAL ERROR in run_scan_and_display: {e}")
        traceback.print_exc()
        _empty_display = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
        _empty_state = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action': '📂', 'Manual_Notes': ''})
        return _empty_display, f"Error during scan: {e}", "", _empty_state, ""


def _text_search_mask(series: pd.Series, terms):
    mask = pd.Series(True, index=series.index)
    s = series.astype(str).str.lower().fillna('')

    def _is_pathy(t: str) -> bool:
        return any(ch in t for ch in ['\\', '/', ':'])

    for term in terms:
        t = term.strip().lower()
        if not t:
            continue
        use_regex = not _is_pathy(t)   # literal for path-like strings
        mask &= s.str.contains(t, na=False, regex=use_regex)
        #mask &= s.str.contains(t, na=False, regex=False)

    return mask


def _persistent_dir():
    try:
        base = os.environ.get("LOCALAPPDATA") or os.path.expanduser("~")
        path = os.path.join(base, "InventoryDashboard")
        os.makedirs(path, exist_ok=True)
        return path
    except Exception:
        return os.getcwd()

RECENT_FOLDERS_FILE = os.path.join(_persistent_dir(), "recent_folders.txt")


def filter_dataframe_display(df_full_from_state: pd.DataFrame, status_filter: str, topic_filter_text: str, filename_filter_text: str) -> pd.DataFrame:
    """
    Filtering improvements:
        * Supports folder:EXCLUDE and IncFolder:INCLUDE as before.
        * ALSO performs text search across File Name + Manual_Notes + Content Hint + Identified Topics (DOCX) + Full Path.
    """
    try:
        empty_df_with_display_columns = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')

        if not isinstance(df_full_from_state, pd.DataFrame) or df_full_from_state.empty:
            return empty_df_with_display_columns

        df_filtered = df_full_from_state.copy()

        # Status filter (unchanged)
        if status_filter and status_filter != "All":
            if 'Status' in df_filtered.columns:
                df_filtered = df_filtered[df_filtered['Status'] == status_filter]
            else:
                return empty_df_with_display_columns.copy()

        # Handle filename filter with folder inclusions/exclusions + cross-column term search
        if filename_filter_text:
            if not set(['File Name', 'Folder Path']).issubset(df_filtered.columns):
                return empty_df_with_display_columns.copy()

            # Split into folder filters and search terms
            folder_excludes = []
            folder_includes = []
            search_terms = []

            for term in filename_filter_text.split(','):
                term = term.strip()
                term_lower = term.lower()
                if term_lower.startswith('folder:'):
                    folder_term = term[7:].strip()
                    if folder_term:
                        folder_excludes.append(folder_term.lower())
                elif term_lower.startswith('incfolder:'):
                    folder_term = term[10:].strip()
                    if folder_term:
                        folder_includes.append(folder_term.lower())
                else:
                    if term:
                        search_terms.append(term)

            
            # Apply folder inclusions (OR logic - keep if matches any term)
            if folder_includes:
                mask = pd.Series(False, index=df_filtered.index)
                fp = df_filtered['Folder Path'].astype(str).str.lower()
                for term in folder_includes:
                    t = term.strip()
                    if t:
                        mask |= fp.str.contains(t, na=False, regex=False)  # literal
                df_filtered = df_filtered[mask]
            
            # Apply folder exclusions (OR logic - exclude if matches any term)
            if folder_excludes:
                fp = df_filtered['Folder Path'].astype(str).str.lower()
                for term in folder_excludes:
                    t = term.strip()
                    if t:
                        df_filtered = df_filtered[~fp.str.contains(t, na=False, regex=False)]  # literal


            
            # Cross-column search (AND across terms)
            if search_terms:
                columns_to_search = [
                    'File Name',
                    'Manual_Notes',
                    'Content Hint',
                    'Identified Topics (DOCX)',
                    'Full Path',
                    'Last Modified'   # <-- add this
                ]
                # Build combined string
                combined = pd.Series('', index=df_filtered.index)
                for col in columns_to_search:
                    if col in df_filtered.columns:
                        combined = combined.str.cat(df_filtered[col].astype(str).fillna(''), sep=' || ')
                mask_all = _text_search_mask(combined, search_terms)
                df_filtered = df_filtered[mask_all]


        # Topic filter (unchanged; AND across terms)
        if topic_filter_text:
            if 'Identified Topics (DOCX)' in df_filtered.columns:
                for term in [t for t in topic_filter_text.split(',') if t.strip()]:
                    df_filtered = df_filtered[
                        df_filtered['Identified Topics (DOCX)'].astype(str).str.lower().str.contains(term.strip().lower(), na=False)
                    ]

        return df_filtered.reindex(columns=DISPLAY_COLUMNS).fillna('')
    except Exception as e:
        print(f"CRITICAL ERROR in filter_dataframe_display: {e}")
        traceback.print_exc()
        return pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')


def save_notes(displayed_df_with_edits: pd.DataFrame, full_df_from_state: pd.DataFrame, xlsx_path: str):
    """
    Update Manual_Notes in the master df state from the displayed (possibly filtered) df,
    and persist to disk with merge protection to avoid note loss.
    """
    try:
        empty_state_df = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action': '📂', 'Manual_Notes': ''})
        if not isinstance(full_df_from_state, pd.DataFrame) or full_df_from_state.empty:
            return full_df_from_state.copy() if isinstance(full_df_from_state, pd.DataFrame) else empty_state_df, "Cannot save: Master data is empty.", "Last saved: Never"

        if not isinstance(displayed_df_with_edits, pd.DataFrame):
            return full_df_from_state.copy(), "No data displayed to save from.", "Last saved: Never"

        updated_full_df = full_df_from_state.copy()  # Has 'Action' column

        if 'Full Path' not in updated_full_df.columns:
            print("CRITICAL: 'Full Path' column missing in full_df_from_state for save_notes.")
            return full_df_from_state.copy(), "Error: 'Full Path' column missing.", "Last saved: Error"

        updated_full_df.set_index('Full Path', inplace=True, drop=False)  # Keep 'Full Path' as column too for reindex

        if not displayed_df_with_edits.empty and \
                'Full Path' in displayed_df_with_edits.columns and \
                'Manual_Notes' in displayed_df_with_edits.columns:
            notes_to_update_df = displayed_df_with_edits.set_index('Full Path')[['Manual_Notes']].fillna('')
            updated_full_df.update(notes_to_update_df)  # Updates 'Manual_Notes' based on index

        updated_full_df.reset_index(drop=True, inplace=True)  # Remove old index, keep 'Full Path' column

        # Ensure final state df has all STATE_COLUMNS and 'Action' is filled
        final_df_for_state = updated_full_df.reindex(columns=STATE_COLUMNS)
        final_df_for_state['Action'] = final_df_for_state['Action'].fillna('📂')
        final_df_for_state['Manual_Notes'] = final_df_for_state['Manual_Notes'].fillna('')

        # default timestamp
        timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        status_msg = "No valid path."
        if xlsx_path and isinstance(xlsx_path, str) and xlsx_path.strip():
            # Save only FIELDNAMES (data columns) to Excel — with merge to preserve existing notes
            if save_inventory_to_xlsx(final_df_for_state[FIELDNAMES].to_dict(orient='records'), xlsx_path):
                status_msg = f"Notes saved successfully to {os.path.basename(xlsx_path)}. (at {timestamp})"
            else:
                status_msg = "Error: Failed to save notes to file."
        else:
            status_msg = "Error: Inventory file path not set."

        return final_df_for_state, status_msg, f"Last saved: {timestamp}"
    except Exception as e:
        print(f"CRITICAL ERROR in save_notes: {e}")
        traceback.print_exc()
        _empty_state = pd.DataFrame(columns=STATE_COLUMNS).fillna({'Action': '📂', 'Manual_Notes': ''})
        return full_df_from_state.copy() if isinstance(full_df_from_state, pd.DataFrame) else _empty_state, f"Error saving notes: {e}", "Last saved: Error"


class ErrorTracker:
    def __init__(self):
        self.error_count = 0
        self.last_error_time = None
        self.max_errors = 3
        self.error_window = 300  # 5 minutes

    def record_error(self):
        now = time.time()
        if self.last_error_time and (now - self.last_error_time) > self.error_window:
            self.error_count = 0
        self.error_count += 1
        self.last_error_time = now
        return self.error_count >= self.max_errors


error_tracker = ErrorTracker()


# =========================
# Gradio UI
# =========================
with gr.Blocks(theme=gr.themes.Soft()) as demo:
    # State holds data + 'Action' column for internal consistency
    # Initial empty state must match this structure
    _initial_empty_state_df = pd.DataFrame(columns=STATE_COLUMNS)
    _initial_empty_state_df['Action'] = '📂'  # Ensure Action column has icon even if 0 rows
    _initial_empty_state_df['Manual_Notes'] = ''

    full_df_state = gr.State(_initial_empty_state_df)
    current_xlsx_path_state = gr.State("")

    # preload recent folders for the dropdown
    _choices, _value = load_recent_folders(initial_choice=START_FOLDER)

    gr.Markdown(f"# File Inventory Dashboard v{VERSION}")
    gr.Markdown("""<p style='font-size: 0.9em; color: #666; margin-bottom: 10px;'>
        <em>Quick Guide:</em> 📂=open folder | Edit 'Manual_Notes' directly | Search terms comma-separated |
        Filters: <code>folder:old</code> (exclude), <code>IncFolder:data</code> (include).
        <br/>Tip: The search box also matches <strong>Manual_Notes, Content Hint, Identified Topics, and Full Path</strong>.
        Removed files with notes are kept so your notes never disappear.</p>""")
    with gr.Row():
        # SINGLE INPUT WINDOW CHANGED: Textbox -> Dropdown with memory (allow typing new values)
        folder_input = gr.Dropdown(
            label="Folder Path (pick recent or type new)",
            choices=_choices,
            value=_value,
            allow_custom_value=True
        )
        scan_button = gr.Button("Scan Folder & Load/Update Inventory", variant="primary")
        shutdown_button = gr.Button("Exit Application", variant="stop")

    status_output = gr.Markdown()
    file_op_status_output = gr.Markdown()
    current_xlsx_display = gr.Textbox(label="Active Inventory File Path", interactive=False)

    gr.Markdown("## Filter and Save")
    with gr.Row():
        status_dropdown = gr.Dropdown(choices=["All", "Active", "Updated", "Added", "Removed (Not Found)"], value="All", label="Filter by Status")
        topic_search = gr.Textbox(label="Search Topics")
        filename_search = gr.Textbox(
            label="Search Name / Notes / Path",
            placeholder="Examples: manuscript, draft | folder:old (exclude) | IncFolder:data (include only)"
        )
    with gr.Row():
        filter_button = gr.Button("Apply Filters")
        save_notes_button = gr.Button("Save Notes", variant="primary")

    gr.Markdown("## Inventory Data (Click '📂' to open folder, edit 'Manual_Notes' column, then click 'Save Notes')")

    # dataframe_output displays DISPLAY_COLUMNS
    _initial_empty_display_df = pd.DataFrame(columns=DISPLAY_COLUMNS).fillna('')
    dataframe_output = gr.DataFrame(
        value=_initial_empty_display_df,
        interactive=True, wrap=True, headers=DISPLAY_COLUMNS,
        datatype=['markdown'] + ['str'] * (len(DISPLAY_COLUMNS) - 1)
    )

    with gr.Row():
        last_save_indicator = gr.Markdown("Last saved: Never")

    # Wrapper: remember folder, then run scan, then update dropdown choices
    def scan_and_remember(folder_path_input):
        # If valid, store to recent
        if folder_path_input and os.path.isdir(str(folder_path_input)):
            add_recent_folder(str(folder_path_input))
            # refresh dropdown choices
            new_choices, _ = load_recent_folders(initial_choice=str(folder_path_input))
            dropdown_update = gr.update(choices=new_choices, value=str(folder_path_input))
        else:
            dropdown_update = gr.update(value=str(folder_path_input) if folder_path_input else None)


        # Original behavior
        
        d, s, x, f, xs = run_scan_and_display(str(folder_path_input).strip() if folder_path_input else "")

        return d, s, x, f, xs, dropdown_update

    scan_button.click(
        fn=scan_and_remember,
        inputs=[folder_input],
        outputs=[dataframe_output, status_output, current_xlsx_display, full_df_state, current_xlsx_path_state, folder_input]
    )

    filter_button.click(
        fn=filter_dataframe_display,
        inputs=[full_df_state, status_dropdown, topic_search, filename_search],
        outputs=[dataframe_output]
    )
    shutdown_button.click(
        fn=shutdown_server,
        inputs=None,
        outputs=gr.Textbox(visible=False)
    )
    dataframe_output.select(
        fn=handle_action_click, inputs=[dataframe_output], outputs=[file_op_status_output]
    )
    save_notes_button.click(
        fn=save_notes,
        inputs=[dataframe_output, full_df_state, current_xlsx_path_state],
        outputs=[full_df_state, file_op_status_output, last_save_indicator]
    )


def setup_shutdown_handler():
    import atexit

    def cleanup():
        print("Performing cleanup before exit...")
        # Save any pending changes
        if hasattr(demo, 'current_xlsx_path_state') and demo.current_xlsx_path_state:
            try:
                create_backup(demo.current_xlsx_path_state)
            except Exception:
                pass

    atexit.register(cleanup)


# Add to main:
if __name__ == "__main__":
    try:
        setup_shutdown_handler()
        logging.info(f"Launching File Inventory Dashboard v{VERSION}")
        logging.info(f"Start folder: {START_FOLDER}")
        demo.launch(
            server_name="127.0.0.1",
            server_port=7860,
            inbrowser=True,
            show_error=True
        )
    except Exception as e:
        logging.error(f"Critical startup error: {e}")
        traceback.print_exc()
